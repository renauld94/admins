#!/usr/bin/env python3
"""
PowerPoint Automation with JNJ Corporate Style
Generates learning materials with Databricks integration for clinical data analysis
Measures impact metrics (relevance to clinical studies, completion rates)
"""

import json
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any


class JNJPowerPointGenerator:
    """Generate PowerPoint presentations with JNJ corporate branding"""
    
    def __init__(self):
        self.output_dir = Path("/home/simon/Learning-Management-System-Academy/staging/powerpoint_generator")
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # JNJ Brand colors
        self.jnj_palette = {
            "primary_red": "#E74C3C",
            "accent_gold": "#F39C12",
            "professional_blue": "#2C3E50",
            "success_green": "#27AE60",
            "light_gray": "#ECF0F1"
        }
        
        # Vietnamese heritage integration
        self.vn_palette = {
            "heritage_red": "#E8423C",
            "heritage_gold": "#C4A73C",
            "heritage_blue": "#1A3A52",
            "heritage_green": "#7BA68F"
        }
    
    def generate_databricks_integration_guide(self) -> str:
        """Generate guide for Databricks clinical data integration"""
        
        guide_content = {
            "title": "Databricks Integration Guide for Vietnamese Course",
            "generated_at": datetime.now().isoformat(),
            "overview": """
This guide enables clinical data analysis and impact measurement 
using Databricks and PySpark for the Vietnamese language learning course.
            """,
            "modules": [
                {
                    "module_id": 101,
                    "module_name": "Foundations of Vietnamese",
                    "clinical_data_points": [
                        "Student engagement metrics (attendance, duration)",
                        "Comprehension assessment scores",
                        "Vocabulary retention rates",
                        "Pronunciation accuracy via audio analysis"
                    ],
                    "databricks_queries": [
                        """
SELECT 
  student_id,
  module_id,
  AVG(engagement_score) as avg_engagement,
  COUNT(DISTINCT day) as days_active,
  AVG(duration_minutes) as avg_session_duration
FROM course_analytics
WHERE module_id = 101
  AND course_id = 10
GROUP BY student_id, module_id
ORDER BY avg_engagement DESC;
                        """
                    ],
                    "impact_metrics": {
                        "primary": "Pronunciation Accuracy (%)",
                        "secondary": ["Vocabulary Retention", "Comprehension", "Engagement"],
                        "target": 85,
                        "current": 78
                    }
                },
                {
                    "module_id": 102,
                    "module_name": "Interactive Communication",
                    "clinical_data_points": [
                        "Conversation practice minutes",
                        "Peer interaction quality",
                        "Response time accuracy",
                        "Vocabulary usage in context"
                    ],
                    "databricks_queries": [
                        """
SELECT 
  student_id,
  COUNT(*) as conversation_sessions,
  AVG(conversation_score) as avg_score,
  MAX(conversation_score) as best_score,
  SUM(practice_minutes) as total_practice_minutes
FROM communication_analytics
WHERE module_id = 102 AND course_id = 10
GROUP BY student_id
HAVING SUM(practice_minutes) > 30
ORDER BY avg_score DESC;
                        """
                    ],
                    "impact_metrics": {
                        "primary": "Conversation Proficiency",
                        "secondary": ["Fluency", "Accuracy", "Confidence"],
                        "target": 88,
                        "current": 82
                    }
                }
            ],
            "pyspark_examples": {
                "setup": """
from pyspark.sql import SparkSession
from pyspark.sql.functions import avg, count, sum as spark_sum, col

# Initialize Spark
spark = SparkSession.builder \\
    .appName("VietnamseCoursAnalytics") \\
    .config("spark.databricks.delta.preview.enabled", "true") \\
    .getOrCreate()

# Load course data
course_analytics = spark.read.table("course_analytics")
student_profiles = spark.read.table("student_profiles")
                """,
                "aggregation": """
# Calculate module-level metrics
module_metrics = course_analytics \\
    .filter((col("course_id") == 10) & (col("status") == "completed")) \\
    .groupBy("module_id") \\
    .agg(
        count("student_id").alias("students_completed"),
        avg("score").alias("avg_score"),
        avg("completion_time_minutes").alias("avg_time"),
        sum("engagement_points").alias("total_engagement")
    ) \\
    .orderBy("module_id")

module_metrics.show()
                """,
                "clinical_relevance": """
# Measure clinical relevance impact
clinical_impact = course_analytics \\
    .join(student_profiles, "student_id") \\
    .groupBy("module_id") \\
    .agg(
        avg(col("clinical_relevance_score")).alias("avg_relevance"),
        count(col("assessment_score") > 80).alias("high_achievers"),
        avg(col("retention_rate")).alias("avg_retention"),
        avg(col("application_to_practice")).alias("practical_application")
    )

clinical_impact.write \\
    .mode("overwrite") \\
    .option("mergeSchema", "true") \\
    .saveAsTable("clinical_impact_metrics")
                """
            },
            "deployment_checklist": [
                "1. Set up Databricks workspace (Community Edition: free for learning)",
                "2. Create Unity Catalog for course data",
                "3. Load course analytics data (student progress, engagement)",
                "4. Configure Moodle API integration with Databricks",
                "5. Build PySpark jobs for nightly metric calculations",
                "6. Create SQL queries for dashboard visualizations",
                "7. Set up automated reports (weekly, monthly)",
                "8. Configure alerts for low engagement/completion",
                "9. Build Databricks SQL warehouse for dashboards",
                "10. Integrate visualizations into Moodle course page"
            ],
            "databricks_advantages": [
                "Real-time analytics on large-scale course data",
                "Clinical data privacy compliance (HIPAA-ready)",
                "Automatic scaling for growing student base",
                "Integration with Moodle via REST APIs",
                "Cost-effective for educational institutions",
                "Support for machine learning workflows",
                "Built-in data governance and audit trails"
            ]
        }
        
        output_file = self.output_dir / "DATABRICKS_INTEGRATION_GUIDE.json"
        with open(output_file, 'w') as f:
            json.dump(guide_content, f, indent=2)
        
        return str(output_file)
    
    def generate_powerpoint_template_structure(self) -> str:
        """Generate PowerPoint template structure with metadata"""
        
        template_structure = {
            "presentation": {
                "title": "Vietnamese Language Learning: Clinical Application",
                "subtitle": "Professional Development Series - Module 1",
                "client": "JNJ Academy",
                "date": datetime.now().strftime("%B %d, %Y"),
                "branding": {
                    "jnj_colors": self.jnj_palette,
                    "vietnamese_heritage": self.vn_palette,
                    "fonts": {
                        "title": "Montserrat Bold",
                        "body": "Open Sans Regular",
                        "accent": "Segoe UI Semibold"
                    }
                },
                "template_pages": []
            },
            "page_templates": {
                "title_slide": {
                    "layout": "Title + Subtitle + Logo",
                    "background": "JNJ Primary Red + Vietnamese Heritage Gold accent",
                    "elements": ["Title", "Subtitle", "Presenter Name", "Date", "JNJ Logo", "Vietnamese Flag"]
                },
                "content_slide": {
                    "layout": "Title + Content + Side Image",
                    "background": "Light Gray",
                    "elements": ["Title", "3-5 bullet points", "Optional: Icon/Image", "Page number"]
                },
                "two_column_slide": {
                    "layout": "Title + Two columns",
                    "background": "White",
                    "elements": ["Title", "Left Column (Text/Image)", "Right Column (Data/Chart)"]
                },
                "statistics_slide": {
                    "layout": "Title + Statistics Grid",
                    "background": "Light Gray with accent borders",
                    "elements": ["Title", "4-6 metric cards with numbers", "Optional: Sparklines"]
                },
                "clinical_impact_slide": {
                    "layout": "Clinical Relevance + Outcomes",
                    "background": "Professional Blue",
                    "elements": ["Title", "Clinical relevance %", "Completion metrics", "Student testimonials"]
                }
            },
            "standard_modules": [
                {
                    "module": 1,
                    "title": "Foundations of Vietnamese",
                    "pages": [
                        {
                            "page": 1,
                            "template": "title_slide",
                            "content": "Module 1: Foundations of Vietnamese"
                        },
                        {
                            "page": 2,
                            "template": "content_slide",
                            "title": "Learning Objectives",
                            "points": [
                                "Master Vietnamese pronunciation fundamentals",
                                "Build foundational vocabulary (200+ essential words)",
                                "Understand cultural context of language use",
                                "Develop listening comprehension skills"
                            ]
                        },
                        {
                            "page": 3,
                            "template": "statistics_slide",
                            "title": "Module 1 Impact Metrics",
                            "metrics": {
                                "students": 450,
                                "completion_rate": 92,
                                "avg_score": 85,
                                "clinical_relevance": 85
                            }
                        }
                    ]
                },
                {
                    "module": 2,
                    "title": "Interactive Communication",
                    "pages": [
                        {
                            "page": 1,
                            "template": "title_slide",
                            "content": "Module 2: Interactive Communication"
                        },
                        {
                            "page": 2,
                            "template": "content_slide",
                            "title": "Clinical Application",
                            "points": [
                                "Patient interaction protocols in Vietnamese",
                                "Medical terminology and expressions",
                                "Active listening and empathy communication",
                                "Handling sensitive healthcare conversations"
                            ]
                        },
                        {
                            "page": 3,
                            "template": "clinical_impact_slide",
                            "title": "Clinical Relevance",
                            "relevance_score": 90,
                            "outcomes": ["Improved patient satisfaction", "+15% communication score", "+22% confidence rating"]
                        }
                    ]
                }
            ],
            "generation_options": {
                "format": "PPTX (Office Open XML)",
                "theme": "JNJ Corporate + Vietnamese Heritage",
                "include_branding": True,
                "include_charts": True,
                "include_speaker_notes": True,
                "generate_pdf_version": True
            },
            "automation_workflow": [
                "1. Read module definitions from MODULE_ID_REFERENCE.md",
                "2. Query Databricks for clinical metrics",
                "3. Generate slides per module (4-6 slides each)",
                "4. Apply JNJ branding + Vietnamese heritage colors",
                "5. Insert D3.js charts as images",
                "6. Add speaker notes for instructors",
                "7. Generate final PPTX file",
                "8. Create PDF version for sharing",
                "9. Log generation metrics to Databricks"
            ]
        }
        
        output_file = self.output_dir / "POWERPOINT_TEMPLATE_STRUCTURE.json"
        with open(output_file, 'w') as f:
            json.dump(template_structure, f, indent=2)
        
        return str(output_file)
    
    def generate_powerpoint_builder_script(self) -> str:
        """Generate Python script for building PowerPoint presentations"""
        
        script = '''#!/usr/bin/env python3
"""
PowerPoint Presentation Builder
Generates PPTX files with JNJ branding and clinical data integration
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    print("✅ python-pptx installed")
except ImportError:
    print("⚠️ python-pptx not installed. Run: pip install python-pptx")

import json
from datetime import datetime
from pathlib import Path


class JNJPowerPointBuilder:
    """Build PPTX presentations with JNJ branding"""
    
    def __init__(self, title: str, module_id: int):
        """Initialize presentation builder"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        self.title = title
        self.module_id = module_id
        
        # JNJ Colors
        self.COLORS = {
            "jnj_red": RGBColor(231, 76, 60),        # #E74C3C
            "jnj_gold": RGBColor(243, 156, 18),      # #F39C12
            "jnj_blue": RGBColor(44, 62, 80),        # #2C3E50
            "vn_red": RGBColor(232, 66, 60),         # #E8423C
            "vn_gold": RGBColor(196, 167, 60),       # #C4A73C
            "vn_blue": RGBColor(26, 58, 82),         # #1A3A52
            "vn_green": RGBColor(123, 166, 143),     # #7BA68F
            "white": RGBColor(255, 255, 255),
            "dark_gray": RGBColor(44, 62, 80)
        }
    
    def add_title_slide(self, title: str, subtitle: str = ""):
        """Add title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS["vn_red"]
        
        # Add accent bar (Vietnamese gold)
        accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.COLORS["vn_gold"]
        accent.line.color.rgb = self.COLORS["vn_gold"]
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.COLORS["white"]
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = self.COLORS["white"]
    
    def add_content_slide(self, title: str, bullet_points: list):
        """Add content slide with bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS["white"]
        
        # Header bar
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
        header.fill.solid()
        header.fill.fore_color.rgb = self.COLORS["vn_blue"]
        header.line.color.rgb = self.COLORS["vn_blue"]
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.COLORS["white"]
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8.5), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = self.COLORS["dark_gray"]
            p.space_before = Pt(12)
    
    def add_metrics_slide(self, title: str, metrics: dict):
        """Add metrics/statistics slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Header bar
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
        header.fill.solid()
        header.fill.fore_color.rgb = self.COLORS["vn_gold"]
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.COLORS["dark_gray"]
        
        # Metrics grid
        cols = 2
        metric_items = list(metrics.items())
        
        for idx, (metric_name, metric_value) in enumerate(metric_items):
            row = idx // cols
            col = idx % cols
            
            x = Inches(1 + col * 4.5)
            y = Inches(1.5 + row * 2.2)
            
            # Metric box
            box = slide.shapes.add_shape(1, x, y, Inches(4), Inches(1.8))
            box.fill.solid()
            box.fill.fore_color.rgb = self.COLORS["vn_green"] if idx % 2 == 0 else self.COLORS["vn_blue"]
            box.line.color.rgb = self.COLORS["vn_gold"]
            box.line.width = Pt(2)
            
            # Metric label
            label_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(3.6), Inches(0.5))
            label_frame = label_box.text_frame
            label_frame.text = metric_name
            label_frame.paragraphs[0].font.size = Pt(12)
            label_frame.paragraphs[0].font.color.rgb = self.COLORS["white"]
            
            # Metric value
            value_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), Inches(3.6), Inches(0.9))
            value_frame = value_box.text_frame
            value_frame.text = str(metric_value)
            value_frame.paragraphs[0].font.size = Pt(36)
            value_frame.paragraphs[0].font.bold = True
            value_frame.paragraphs[0].font.color.rgb = self.COLORS["white"]
    
    def save(self, filename: str) -> str:
        """Save presentation"""
        output_path = Path(f"/home/simon/Learning-Management-System-Academy/staging/powerpoint_generator/{filename}")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        self.prs.save(str(output_path))
        return str(output_path)


def create_sample_presentation():
    """Create sample presentation for Module 1"""
    prs_builder = JNJPowerPointBuilder(
        "Vietnamese Language: Foundations",
        module_id=101
    )
    
    # Slide 1: Title
    prs_builder.add_title_slide(
        "Foundations of Vietnamese",
        "Professional Development Series - Module 1"
    )
    
    # Slide 2: Learning Objectives
    prs_builder.add_content_slide(
        "Learning Objectives",
        [
            "Master Vietnamese pronunciation fundamentals",
            "Build foundational vocabulary (200+ essential words)",
            "Understand cultural context of language use",
            "Develop listening comprehension skills"
        ]
    )
    
    # Slide 3: Module Metrics
    prs_builder.add_metrics_slide(
        "Module 1 Impact Metrics",
        {
            "Students": "450",
            "Completion": "92%",
            "Avg Score": "85%",
            "Relevance": "85%"
        }
    )
    
    # Slide 4: Content
    prs_builder.add_content_slide(
        "Key Content Areas",
        [
            "Pronunciation: Tones and phonetics",
            "Vocabulary: Greetings, numbers, basic phrases",
            "Grammar: Word order and particles",
            "Cultural: Etiquette and social norms"
        ]
    )
    
    output = prs_builder.save("Module_1_Foundations_Vietnamese.pptx")
    return output


if __name__ == "__main__":
    import sys
    
    # Check if python-pptx is installed
    try:
        from pptx import Presentation
        output_file = create_sample_presentation()
        print(f"✅ Generated: {output_file}")
    except ImportError:
        print("⚠️ python-pptx not installed")
        print("   To generate PowerPoint files, run:")
        print("   pip install python-pptx")
        print("\\nAlternatively, use Office Open XML API or LibreOffice Calc")
        sys.exit(1)
'''
        
        output_file = self.output_dir / "powerpoint_builder.py"
        with open(output_file, 'w') as f:
            f.write(script)
        
        return str(output_file)
    
    def generate_all_resources(self) -> Dict[str, str]:
        """Generate all PowerPoint automation resources"""
        
        print("\n" + "="*80)
        print("📊 PowerPoint Automation with JNJ Branding & Databricks Integration")
        print("="*80)
        
        results = {}
        
        print("\n[1/3] Generating Databricks integration guide...")
        results["databricks_guide"] = self.generate_databricks_integration_guide()
        print(f"✅ Databricks Guide: {results['databricks_guide']}")
        
        print("\n[2/3] Generating PowerPoint template structure...")
        results["template_structure"] = self.generate_powerpoint_template_structure()
        print(f"✅ Template Structure: {results['template_structure']}")
        
        print("\n[3/3] Generating PowerPoint builder script...")
        results["builder_script"] = self.generate_powerpoint_builder_script()
        print(f"✅ Builder Script: {results['builder_script']}")
        
        print("\n" + "="*80)
        print("✨ PowerPoint automation resources generated successfully!")
        print("="*80 + "\n")
        
        return results


def main():
    generator = JNJPowerPointGenerator()
    resources = generator.generate_all_resources()
    
    print("\n📁 PowerPoint Automation Resources:")
    for name, path in resources.items():
        print(f"  • {name}: {path}")


if __name__ == "__main__":
    main()
