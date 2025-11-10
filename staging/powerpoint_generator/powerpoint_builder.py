#!/usr/bin/env python3
"""
PowerPoint Presentation Builder
Generates PPTX files with JNJ branding and clinical data integration
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    print("✅ python-pptx installed")
except ImportError:
    print("⚠️ python-pptx not installed. Run: pip install python-pptx")

import json
from datetime import datetime
from pathlib import Path


class JNJPowerPointBuilder:
    """Build PPTX presentations with JNJ branding"""
    
    def __init__(self, title: str, module_id: int):
        """Initialize presentation builder"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        self.title = title
        self.module_id = module_id
        
        # JNJ Colors
        self.COLORS = {
            "jnj_red": RGBColor(231, 76, 60),        # #E74C3C
            "jnj_gold": RGBColor(243, 156, 18),      # #F39C12
            "jnj_blue": RGBColor(44, 62, 80),        # #2C3E50
            "vn_red": RGBColor(232, 66, 60),         # #E8423C
            "vn_gold": RGBColor(196, 167, 60),       # #C4A73C
            "vn_blue": RGBColor(26, 58, 82),         # #1A3A52
            "vn_green": RGBColor(123, 166, 143),     # #7BA68F
            "white": RGBColor(255, 255, 255),
            "dark_gray": RGBColor(44, 62, 80)
        }
    
    def add_title_slide(self, title: str, subtitle: str = ""):
        """Add title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS["vn_red"]
        
        # Add accent bar (Vietnamese gold)
        accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.COLORS["vn_gold"]
        accent.line.color.rgb = self.COLORS["vn_gold"]
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.COLORS["white"]
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = self.COLORS["white"]
    
    def add_content_slide(self, title: str, bullet_points: list):
        """Add content slide with bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS["white"]
        
        # Header bar
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
        header.fill.solid()
        header.fill.fore_color.rgb = self.COLORS["vn_blue"]
        header.line.color.rgb = self.COLORS["vn_blue"]
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.COLORS["white"]
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8.5), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = self.COLORS["dark_gray"]
            p.space_before = Pt(12)
    
    def add_metrics_slide(self, title: str, metrics: dict):
        """Add metrics/statistics slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Header bar
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.7))
        header.fill.solid()
        header.fill.fore_color.rgb = self.COLORS["vn_gold"]
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.COLORS["dark_gray"]
        
        # Metrics grid
        cols = 2
        metric_items = list(metrics.items())
        
        for idx, (metric_name, metric_value) in enumerate(metric_items):
            row = idx // cols
            col = idx % cols
            
            x = Inches(1 + col * 4.5)
            y = Inches(1.5 + row * 2.2)
            
            # Metric box
            box = slide.shapes.add_shape(1, x, y, Inches(4), Inches(1.8))
            box.fill.solid()
            box.fill.fore_color.rgb = self.COLORS["vn_green"] if idx % 2 == 0 else self.COLORS["vn_blue"]
            box.line.color.rgb = self.COLORS["vn_gold"]
            box.line.width = Pt(2)
            
            # Metric label
            label_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(3.6), Inches(0.5))
            label_frame = label_box.text_frame
            label_frame.text = metric_name
            label_frame.paragraphs[0].font.size = Pt(12)
            label_frame.paragraphs[0].font.color.rgb = self.COLORS["white"]
            
            # Metric value
            value_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.8), Inches(3.6), Inches(0.9))
            value_frame = value_box.text_frame
            value_frame.text = str(metric_value)
            value_frame.paragraphs[0].font.size = Pt(36)
            value_frame.paragraphs[0].font.bold = True
            value_frame.paragraphs[0].font.color.rgb = self.COLORS["white"]
    
    def save(self, filename: str) -> str:
        """Save presentation"""
        output_path = Path(f"/home/simon/Learning-Management-System-Academy/staging/powerpoint_generator/{filename}")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        self.prs.save(str(output_path))
        return str(output_path)


def create_sample_presentation():
    """Create sample presentation for Module 1"""
    prs_builder = JNJPowerPointBuilder(
        "Vietnamese Language: Foundations",
        module_id=101
    )
    
    # Slide 1: Title
    prs_builder.add_title_slide(
        "Foundations of Vietnamese",
        "Professional Development Series - Module 1"
    )
    
    # Slide 2: Learning Objectives
    prs_builder.add_content_slide(
        "Learning Objectives",
        [
            "Master Vietnamese pronunciation fundamentals",
            "Build foundational vocabulary (200+ essential words)",
            "Understand cultural context of language use",
            "Develop listening comprehension skills"
        ]
    )
    
    # Slide 3: Module Metrics
    prs_builder.add_metrics_slide(
        "Module 1 Impact Metrics",
        {
            "Students": "450",
            "Completion": "92%",
            "Avg Score": "85%",
            "Relevance": "85%"
        }
    )
    
    # Slide 4: Content
    prs_builder.add_content_slide(
        "Key Content Areas",
        [
            "Pronunciation: Tones and phonetics",
            "Vocabulary: Greetings, numbers, basic phrases",
            "Grammar: Word order and particles",
            "Cultural: Etiquette and social norms"
        ]
    )
    
    output = prs_builder.save("Module_1_Foundations_Vietnamese.pptx")
    return output


if __name__ == "__main__":
    import sys
    
    # Check if python-pptx is installed
    try:
        from pptx import Presentation
        output_file = create_sample_presentation()
        print(f"✅ Generated: {output_file}")
    except ImportError:
        print("⚠️ python-pptx not installed")
        print("   To generate PowerPoint files, run:")
        print("   pip install python-pptx")
        print("\nAlternatively, use Office Open XML API or LibreOffice Calc")
        sys.exit(1)
