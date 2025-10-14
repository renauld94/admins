from pptx import Presentation

input_path = "/home/simon/Downloads/Empowering-Clinical-Programming-with-Python-Excellence (1).pptx"
output_path = "/home/simon/Downloads/Empowering-Clinical-Programming-with-Python-Excellence-Transformed.pptx"
prs = Presentation(input_path)

# Remove footer shapes from all slides
for slide in prs.slides:
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            text = shape.text_frame.text.lower()
            if "page" in text or "footer" in text or "confidential" in text:
                slide.shapes._spTree.remove(shape._element)
modules = [
    {
        "title": "Module 1: Core Python Foundations (2h 00m)",
        "sessions": [
            ("1.1", "Basics", "Essential foundation (variables, expressions, types)"),
            ("1.2", "Modules and Packages", "Supports modularity, used in PySpark imports"),
            ("1.3", "Data Structures", "Core for manipulating lists, dicts, comprehensions"),
            ("1.4", "Advanced Data Structures", "Important for handling nested schemas in DataFrames"),
            ("1.5", "Conditions and Loops", "Critical for control flow and PySpark logic"),
            ("1.6", "Functions", "Includes lambda, map, filter – core concepts in PySpark"),
            ("1.7", "Dates and Times", "Important for timestamp data handling"),
            ("1.8", "Regular Expressions", "Useful for data cleaning and parsing text"),
            ("1.9", "Classes", "Builds understanding of structure and abstraction"),
            ("1.10", "Decorators", "Supports advanced functions, relevant for modular pipelines"),
            ("1.11", "Virtual Environments", "Good practice for environment isolation"),
            ("1.12", "Exception Handling", "Needed for building robust PySpark data pipelines"),
            ("1.13", "Context Managers and File I/O", "Essential for reading/writing files and managing resources"),
            ("1.14", "Iterators and Generators", "Maps directly to lazy evaluation in PySpark"),
            ("1.15", "String Processing", "Important for cleaning and manipulating string data"),
            ("1.16", "APIs and JSON", "Common for ingesting data from APIs and parsing input files"),
        ]
    },
    {
        "title": "Module 2: PySpark Fundamentals (2h 15m)",
        "sessions": [
            ("2.1", "Welcome & Introduction", "Orientation to distributed computing"),
            ("2.2", "Setting Up PySpark", "Environment setup, SparkSession"),
            ("2.3", "SparkSession and SparkContext", "Core Spark concepts"),
            ("2.4", "RDD Fundamentals", "Low-level distributed data structures"),
            ("2.5", "DataFrames Basics", "Main data structure for analytics"),
            ("2.6", "DataFrame Operations", "Transformations and actions"),
            ("2.7", "PySpark SQL", "SQL queries on big data"),
            ("2.8", "User Defined Functions (UDFs)", "Custom logic in Spark"),
            ("2.9", "Working with Dates & Timestamps", "Time-based analytics"),
            ("2.10", "Window Functions", "Advanced analytics"),
            ("2.11", "Reading and Writing Data", "Data I/O"),
            ("2.12", "Performance Tuning Basics", "Optimizing Spark jobs"),
            ("2.13", "Error Handling and Debugging", "Building robust pipelines"),
            ("2.14", "Practical Exercises", "Hands-on practice"),
        ]
    },
    {
        "title": "Module 3: PySpark Advanced (1h 30m)",
        "sessions": [
            ("3.1", "Complex Operations", "Multi-step transformations"),
            ("3.2", "Joins and Aggregations", "Combining and summarizing data"),
            ("3.3", "Performance Optimization", "Advanced tuning, caching, partitioning"),
            ("3.4", "Advanced Functions", "Window, array, map, struct functions"),
            ("3.5", "Error Handling & Debugging", "Troubleshooting Spark jobs"),
            ("3.6", "Capstone Project", "Real-world scenario, end-to-end pipeline"),
        ]
    },
    {
        "title": "Module 4: Git & Bitbucket for Data Teams (1h 00m)",
        "sessions": [
            ("4.1", "Welcome & Introduction", "Version control basics"),
            ("4.2", "Introduction to Version Control", "Git fundamentals"),
            ("4.3", "Git Basics and Workflow", "Branching, merging, workflow"),
            ("4.4", "Working with Bitbucket", "Cloud repository management"),
            ("4.5", "Branching and Merging", "Collaboration"),
            ("4.6", "Collaboration and Pull Requests", "Team workflows"),
            ("4.7", "Resolving Conflicts", "Conflict management"),
            ("4.8", "Best Practices", "Professional standards"),
        ]
    },
    {
        "title": "Module 5: Databricks Platform (1h 00m)",
        "sessions": [
            ("5.1", "Welcome & Introduction", "Platform orientation"),
            ("5.2", "Introduction to Databricks", "Key features, workspace"),
            ("5.3", "Setting Up a Workspace", "Environment setup"),
            ("5.4", "Running Notebooks", "Interactive development"),
            ("5.5", "Integrating with PySpark", "Unified analytics"),
            ("5.6", "Databricks Utilities", "dbutils, file management"),
            ("5.7", "Job Scheduling", "Automation"),
            ("5.8", "Best Practices", "Efficient workflows"),
        ]
    },
    {
        "title": "Module 6: Advanced Data Engineering (1h 15m)",
        "sessions": [
            ("6.1", "Welcome & Introduction", "Advanced topics overview"),
            ("6.2", "Advanced Python Concepts", "Generators, decorators, context managers"),
            ("6.3", "Advanced PySpark Techniques", "Broadcast joins, partitioning, custom UDFs"),
            ("6.4", "Performance Optimization", "Profiling, tuning, scaling"),
            ("6.5", "Real-world Data Pipelines", "End-to-end pipeline design"),
            ("6.6", "Testing and Debugging", "Unit tests, error handling"),
            ("6.7", "Deployment Strategies", "CI/CD, packaging, cloud deployment"),
            ("6.8", "Capstone Project", "Comprehensive project, real healthcare data"),
        ]
    }
]

# Print available layouts for debugging
for idx, layout in enumerate(prs.slide_layouts):
    print(f"Layout {idx}: {layout.name}")

# Use the first available layout if index 1 is not present
if len(prs.slide_layouts) > 1:
    content_layout = prs.slide_layouts[1]
else:
    content_layout = prs.slide_layouts[0]

for module in modules:
    slide = prs.slides.add_slide(content_layout)
    # Try to set the title, or add a textbox if not present
    if slide.shapes.title:
        slide.shapes.title.text = module["title"]
    else:
        left = top = width = height = prs.slide_width // 20, prs.slide_height // 20, prs.slide_width * 9 // 10, prs.slide_height // 8
        textbox = slide.shapes.add_textbox(*left, *top, *width, *height)
        textbox.text = module["title"]
    # Try to set the content, or add a textbox if not present
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Sessions:"
    else:
        left = prs.slide_width // 20
        top = prs.slide_height // 6
        width = prs.slide_width * 9 // 10
        height = prs.slide_height * 3 // 5
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text = "Sessions:"

    for session in module["sessions"]:
        session_slide = prs.slides.add_slide(content_layout)
        if session_slide.shapes.title:
            session_slide.shapes.title.text = f"{module['title']} - {session[0]}"
        else:
            left = top = width = height = prs.slide_width // 20, prs.slide_height // 20, prs.slide_width * 9 // 10, prs.slide_height // 8
            textbox = session_slide.shapes.add_textbox(*left, *top, *width, *height)
            textbox.text = f"{module['title']} - {session[0]}"
        if len(session_slide.placeholders) > 1:
            session_slide.placeholders[1].text = f"Title: {session[1]}\nRelevance: {session[2]}"
        else:
            left = prs.slide_width // 20
            top = prs.slide_height // 6
            width = prs.slide_width * 9 // 10
            height = prs.slide_height * 3 // 5
            textbox = session_slide.shapes.add_textbox(left, top, width, height)
            textbox.text = f"Title: {session[1]}\nRelevance: {session[2]}"

prs.save(output_path)
print(f"Transformed presentation saved as: {output_path}")